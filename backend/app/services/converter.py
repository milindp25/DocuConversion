"""
PDF conversion service.

Handles the core logic for converting between PDF and other document
formats. Uses PyMuPDF for PDF reading, python-docx for Word output,
Pillow for image processing, and ReportLab for PDF generation.
"""

import io
import logging
import textwrap
import zipfile
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, Inches, RGBColor
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches, Emu
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from app.core.exceptions import ConversionError

logger = logging.getLogger(__name__)


class ConversionService:
    """Service for converting documents between PDF and other formats.

    Each conversion method reads the source file, processes it using
    the appropriate library, and returns the path to the output file.
    All methods are designed to be called from the API layer.
    """

    @staticmethod
    async def pdf_to_text(input_path: Path) -> str:
        """Extract all text content from a PDF document.

        Uses PyMuPDF's text extraction which preserves basic structure
        like paragraphs and line breaks.

        Args:
            input_path: Path to the source PDF file.

        Returns:
            Extracted text content as a string.

        Raises:
            ConversionError: If the PDF cannot be read or processed.
        """
        try:
            with fitz.open(str(input_path)) as doc:
                text_parts: list[str] = []

                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text_parts.append(page.get_text())

                return "\n\n".join(text_parts)

        except Exception as e:
            logger.exception("PDF to text extraction failed:")
            raise ConversionError(
                "Failed to extract text from the PDF. "
                "The file may be corrupted or contain only scanned images. "
                "Try using OCR for scanned documents."
            ) from e

    @staticmethod
    async def pdf_to_text_file(input_path: Path, output_path: Path) -> Path:
        """Extract text from a PDF and save it to a UTF-8 text file.

        Delegates to ``pdf_to_text`` for extraction, then writes the
        resulting string to *output_path* as a plain ``.txt`` file.

        Args:
            input_path: Path to the source PDF file.
            output_path: Where to save the generated text file.

        Returns:
            Path to the created text file.

        Raises:
            ConversionError: If the PDF cannot be read or text file creation fails.
        """
        text = await ConversionService.pdf_to_text(input_path)
        try:
            output_path.write_text(text, encoding="utf-8")
            logger.info(
                "PDF to text file conversion complete: %s -> %s",
                input_path.name,
                output_path.name,
            )
            return output_path
        except Exception as e:
            logger.exception("PDF to text file write failed:")
            raise ConversionError(
                "Failed to write extracted text to file."
            ) from e

    @staticmethod
    def _detect_font_sizes(pdf_doc: fitz.Document) -> dict[str, float]:
        """Analyze the PDF to determine font size thresholds for headings.

        Scans all spans across all pages, counts font size frequency,
        and classifies sizes into body, heading, and title categories.

        Returns:
            Dict with 'body', 'heading', and 'title' size thresholds.
        """
        size_counts: dict[float, int] = {}
        for page in pdf_doc:
            data = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            for block in data.get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        size = round(span["size"], 1)
                        size_counts[size] = size_counts.get(size, 0) + 1

        if not size_counts:
            return {"body": 10.0, "heading": 12.0, "title": 16.0}

        # The most frequent size is body text
        body_size = max(size_counts, key=size_counts.get)
        return {
            "body": body_size,
            "heading": body_size * 1.15,
            "title": body_size * 1.5,
        }

    @staticmethod
    def _fix_encoding(text: str) -> str:
        """Fix common Unicode replacement characters from PDF extraction.

        Some PDF fonts use non-standard character maps that cause PyMuPDF
        to emit U+FFFD (replacement character) for en-dashes, em-dashes,
        bullets, and other special characters.  This method replaces the
        most common occurrences based on surrounding context.

        Args:
            text: Raw text from PyMuPDF span extraction.

        Returns:
            Text with replacement characters fixed.
        """
        import re

        en_dash = "\u2013"  # –
        bullet = "\u2022"   # •

        # U+FFFD between digits → en-dash (e.g. "900–6017")
        text = re.sub(
            r"(\d)\ufffd(\d)",
            lambda m: m.group(1) + en_dash + m.group(2),
            text,
        )
        # U+FFFD between spaces/words → en-dash (e.g. "May 2025 – Present")
        text = re.sub(
            r"\s\ufffd\s",
            lambda m: " " + en_dash + " ",
            text,
        )
        # U+FFFD at start of text → bullet
        if text.startswith("\ufffd"):
            text = bullet + text[1:]
        # Any remaining U+FFFD → en-dash
        text = text.replace("\ufffd", en_dash)
        return text

    @staticmethod
    def _should_split_line(prev_line: dict, curr_line: dict) -> bool:
        """Determine if two consecutive lines should be separate paragraphs.

        Lines within the same PyMuPDF block may belong to different
        logical paragraphs (e.g., a company name followed by a job title).
        This method detects formatting changes that indicate a paragraph
        break.

        Args:
            prev_line: The preceding line dict from ``get_text("dict")``.
            curr_line: The current line dict.

        Returns:
            True if a paragraph break should be inserted.
        """
        prev_spans = prev_line.get("spans", [])
        curr_spans = curr_line.get("spans", [])
        if not prev_spans or not curr_spans:
            return False

        prev_text = "".join(s["text"] for s in prev_spans).strip()
        curr_text = "".join(s["text"] for s in curr_spans).strip()

        # Bullet lines always start a new paragraph
        if curr_text and curr_text[0] in ("\u2022", "\u2023", "\u25e6", "•", "\ufffd"):
            return True

        # Determine dominant formatting of each line
        prev_bold = any(
            s["flags"] & (1 << 4) for s in prev_spans if s["text"].strip()
        )
        prev_italic = any(
            s["flags"] & (1 << 1) for s in prev_spans if s["text"].strip()
        )
        curr_bold = any(
            s["flags"] & (1 << 4) for s in curr_spans if s["text"].strip()
        )
        curr_italic = any(
            s["flags"] & (1 << 1) for s in curr_spans if s["text"].strip()
        )

        # Italic line always starts a new paragraph (job titles, etc.)
        if curr_italic and not prev_italic:
            return True

        # Line after italic also starts new paragraph
        if prev_italic and not curr_italic:
            return True

        # Bold line following non-bold starts new paragraph
        if curr_bold and not prev_bold:
            return True

        return False

    @staticmethod
    def _merge_line_spans(spans: list[dict]) -> list[dict]:
        """Merge adjacent spans with identical formatting into one span.

        PyMuPDF often splits a single word or phrase into many tiny
        spans even when their formatting is identical.  Merging them
        produces cleaner Word output with fewer runs per paragraph.

        Args:
            spans: Raw span dicts from PyMuPDF ``get_text("dict")``.

        Returns:
            List of merged span dicts with concatenated text.
        """
        if not spans:
            return []

        merged: list[dict] = []
        current = dict(spans[0])

        for span in spans[1:]:
            same_font = span["font"] == current["font"]
            same_size = abs(span["size"] - current["size"]) < 0.5
            same_flags = span["flags"] == current["flags"]
            same_color = span.get("color") == current.get("color")

            if same_font and same_size and same_flags and same_color:
                current["text"] += span["text"]
            else:
                merged.append(current)
                current = dict(span)

        merged.append(current)
        return merged

    @staticmethod
    async def pdf_to_word(input_path: Path, output_path: Path) -> Path:
        """Convert a PDF document to Word (.docx) format.

        Uses PyMuPDF's ``get_text("dict")`` to extract rich formatting
        data (font name, size, bold, italic, color, position) and
        recreates it in the Word document with proper:
          - Headings detected by font size
          - Bold and italic preserved per span
          - Bullet points from ``•`` / ``-`` prefixes
          - Right-aligned text (e.g. dates) via tab stops
          - Font size proportionally mapped

        Args:
            input_path: Path to the source PDF file.
            output_path: Where to save the generated Word document.

        Returns:
            Path to the created Word document.

        Raises:
            ConversionError: If the PDF cannot be read or Word creation fails.
        """
        try:
            doc_word = Document()

            # Set default font
            style = doc_word.styles["Normal"]
            font = style.font
            font.name = "Calibri"
            font.size = Pt(10)

            with fitz.open(str(input_path)) as pdf_doc:
                total_pages = len(pdf_doc)
                thresholds = ConversionService._detect_font_sizes(pdf_doc)

                for page_num in range(total_pages):
                    page = pdf_doc[page_num]
                    page_width = page.rect.width
                    page_data = page.get_text(
                        "dict", flags=fitz.TEXT_PRESERVE_WHITESPACE
                    )

                    for block in page_data.get("blocks", []):
                        if block.get("type") != 0:
                            continue

                        lines = block.get("lines", [])
                        if not lines:
                            continue

                        # Fix encoding in all spans upfront
                        for line in lines:
                            for span in line.get("spans", []):
                                span["text"] = ConversionService._fix_encoding(
                                    span["text"]
                                )

                        # Split block lines into logical paragraph groups.
                        # Each group becomes one Word paragraph.
                        para_groups: list[list[dict]] = []
                        current_group: list[dict] = [lines[0]]

                        for i in range(1, len(lines)):
                            if ConversionService._should_split_line(
                                lines[i - 1], lines[i]
                            ):
                                para_groups.append(current_group)
                                current_group = [lines[i]]
                            else:
                                current_group.append(lines[i])
                        para_groups.append(current_group)

                        # Process each paragraph group
                        for group_lines in para_groups:

                            # Collect left and right spans
                            left_spans: list[dict] = []
                            right_spans: list[dict] = []

                            for line in group_lines:
                                line_x0 = line["bbox"][0]
                                line_spans = line.get("spans", [])

                                is_right = line_x0 > page_width * 0.65

                                if is_right and left_spans:
                                    right_spans.extend(line_spans)
                                else:
                                    left_spans.extend(line_spans)

                            # Merge tiny spans with identical formatting
                            left_merged = ConversionService._merge_line_spans(
                                left_spans
                            )
                            right_merged = ConversionService._merge_line_spans(
                                right_spans
                            )

                            if not left_merged and not right_merged:
                                continue

                            # Combine all text to determine paragraph type
                            full_text = "".join(
                                s["text"] for s in left_merged
                            ).strip()
                            if not full_text:
                                continue

                            # Skip page number footers (standalone digits
                            # near the bottom of the page)
                            if full_text.isdigit() and len(full_text) <= 3:
                                line_y = group_lines[-1]["bbox"][3]
                                if line_y > page.rect.height * 0.9:
                                    continue

                            # Determine the dominant font size
                            sizes = [
                                s["size"]
                                for s in left_merged
                                if s["text"].strip()
                            ]
                            dominant_size = (
                                max(set(sizes), key=sizes.count)
                                if sizes
                                else thresholds["body"]
                            )

                            # Detect bullet points
                            is_bullet = full_text.startswith(
                                ("\u2022", "\u2023", "\u25e6", "•")
                            )
                            if is_bullet:
                                for s in left_merged:
                                    txt = s["text"].lstrip()
                                    if txt and txt[0] in (
                                        "\u2022", "\u2023", "\u25e6", "•",
                                    ):
                                        s["text"] = txt[1:]
                                        break

                            # Detect heading level
                            is_title = dominant_size >= thresholds["title"]
                            is_heading = (
                                dominant_size >= thresholds["heading"]
                                and not is_title
                            )

                            # Detect centered text
                            avg_x0 = sum(
                                l["bbox"][0] for l in group_lines
                            ) / len(group_lines)
                            avg_x1 = sum(
                                l["bbox"][2] for l in group_lines
                            ) / len(group_lines)
                            left_margin = avg_x0
                            right_margin = page_width - avg_x1
                            is_centered = (
                                abs(left_margin - right_margin)
                                < page_width * 0.15
                                and left_margin > page_width * 0.1
                            )

                            # Create the Word paragraph
                            if is_bullet:
                                para = doc_word.add_paragraph(
                                    style="List Bullet"
                                )
                            else:
                                para = doc_word.add_paragraph()

                            if is_centered:
                                para.alignment = WD_ALIGN_PARAGRAPH.CENTER

                            # Add left-side runs with formatting
                            for span in left_merged:
                                text = span["text"]
                                if not text:
                                    continue

                                run = para.add_run(text)
                                span_size = span["size"]
                                flags = span["flags"]

                                if flags & (1 << 4):
                                    run.bold = True
                                if flags & (1 << 1):
                                    run.italic = True

                                if is_title:
                                    run.font.size = Pt(16)
                                elif is_heading:
                                    run.font.size = Pt(13)
                                else:
                                    ratio = span_size / thresholds["body"]
                                    mapped = max(8, min(14, 10 * ratio))
                                    run.font.size = Pt(round(mapped, 1))

                                color = span.get("color", 0)
                                if color and color != 0:
                                    r_val = (color >> 16) & 0xFF
                                    g_val = (color >> 8) & 0xFF
                                    b_val = color & 0xFF
                                    if (r_val, g_val, b_val) != (0, 0, 0):
                                        run.font.color.rgb = RGBColor(
                                            r_val, g_val, b_val
                                        )

                            # Add right-aligned text (dates) via tab stop
                            if right_merged:
                                right_text = "".join(
                                    s["text"] for s in right_merged
                                ).strip()
                                if right_text:
                                    run = para.add_run("\t" + right_text)
                                    run.font.size = Pt(10)

                                    from docx.oxml import OxmlElement
                                    from docx.oxml.ns import qn
                                    pPr = para._p.get_or_add_pPr()
                                    tabs_el = OxmlElement("w:tabs")
                                    tab_el = OxmlElement("w:tab")
                                    tab_el.set(qn("w:val"), "right")
                                    tab_el.set(qn("w:pos"), "9360")
                                    tabs_el.append(tab_el)
                                    pPr.append(tabs_el)

                    # Page break between pages (except the last)
                    if page_num < total_pages - 1:
                        doc_word.add_page_break()

            doc_word.save(str(output_path))

            logger.info(
                "PDF to Word conversion complete: %d pages -> %s",
                total_pages,
                output_path.name,
            )
            return output_path

        except Exception as e:
            logger.exception("PDF to Word conversion failed:")
            raise ConversionError(
                "Failed to convert PDF to Word format. "
                "The file may be corrupted or contain unsupported content."
            ) from e

    @staticmethod
    async def pdf_to_images(
        input_path: Path, output_dir: Path, dpi: int = 200
    ) -> list[Path]:
        """Convert each page of a PDF to a PNG image.

        Renders each page at the specified DPI for high-quality output.

        Args:
            input_path: Path to the source PDF file.
            output_dir: Directory to save the output images.
            dpi: Resolution for rendering (dots per inch). Default 200.

        Returns:
            List of paths to the generated image files.

        Raises:
            ConversionError: If rendering fails.
        """
        try:
            with fitz.open(str(input_path)) as doc:
                if len(doc) > 300:
                    raise ConversionError(
                        "PDF has too many pages (max 300)."
                    )
                output_paths: list[Path] = []
                zoom = dpi / 72  # PyMuPDF uses 72 DPI as base

                for page_num in range(len(doc)):
                    page = doc[page_num]
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)

                    output_path = output_dir / f"page_{page_num + 1}.png"
                    pix.save(str(output_path))
                    output_paths.append(output_path)

                logger.info(
                    "PDF to images conversion complete: %d pages rendered at %d DPI",
                    len(output_paths),
                    dpi,
                )
                return output_paths

        except Exception as e:
            logger.exception("PDF to image conversion failed:")
            raise ConversionError(
                "Failed to convert PDF pages to images. The file may be corrupted."
            ) from e

    @staticmethod
    async def pdf_to_images_zip(
        input_path: Path, output_path: Path, dpi: int = 200
    ) -> Path:
        """Convert each page of a PDF to PNG images and bundle as a ZIP.

        Convenience method that creates a ZIP archive containing all
        page images, suitable for single-file download.

        Args:
            input_path: Path to the source PDF file.
            output_path: Path for the output ZIP file.
            dpi: Resolution for rendering. Default 200.

        Returns:
            Path to the ZIP archive containing all page images.

        Raises:
            ConversionError: If rendering or archiving fails.
        """
        try:
            with fitz.open(str(input_path)) as doc:
                if len(doc) > 300:
                    raise ConversionError(
                        "PDF has too many pages (max 300)."
                    )
                zoom = dpi / 72

                with zipfile.ZipFile(str(output_path), "w", zipfile.ZIP_DEFLATED) as zf:
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        mat = fitz.Matrix(zoom, zoom)
                        pix = page.get_pixmap(matrix=mat)

                        img_name = f"page_{page_num + 1}.png"
                        zf.writestr(img_name, pix.tobytes("png"))

                logger.info(
                    "PDF to images ZIP complete: %d pages at %d DPI -> %s",
                    len(doc),
                    dpi,
                    output_path.name,
                )
                return output_path

        except Exception as e:
            logger.exception("PDF to images ZIP failed:")
            raise ConversionError(
                "Failed to convert PDF pages to images. The file may be corrupted."
            ) from e

    @staticmethod
    async def image_to_pdf(input_path: Path, output_path: Path) -> Path:
        """Convert an image file to a PDF document.

        Opens the image with Pillow, determines its dimensions, and
        creates a PDF page sized to match using ReportLab. The image
        is fitted to the page while preserving its aspect ratio.

        Args:
            input_path: Path to the source image file.
            output_path: Where to save the generated PDF.

        Returns:
            Path to the created PDF file.

        Raises:
            ConversionError: If the image cannot be read or PDF creation fails.
        """
        try:
            with Image.open(str(input_path)) as img:
                # Convert RGBA to RGB if necessary (PDF doesn't support alpha)
                if img.mode == "RGBA":
                    background = Image.new("RGB", img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[3])
                    img = background
                elif img.mode not in ("RGB", "L"):
                    img = img.convert("RGB")

                img_width, img_height = img.size

                # Create a page sized to the image (in points, 72 DPI)
                # Cap at reasonable page size, scale down if needed
                max_dimension = 14400  # 200 inches at 72 DPI
                scale = min(1.0, max_dimension / max(img_width, img_height))
                page_width = img_width * scale
                page_height = img_height * scale

                # Save image to a bytes buffer for ReportLab
                img_buffer = io.BytesIO()
                img.save(img_buffer, format="PNG")
                img_buffer.seek(0)

            c = canvas.Canvas(str(output_path), pagesize=(page_width, page_height))
            c.drawImage(
                ImageReader(img_buffer),
                0,
                0,
                width=page_width,
                height=page_height,
            )
            c.save()

            logger.info(
                "Image to PDF conversion complete: %dx%d -> %s",
                img_width,
                img_height,
                output_path.name,
            )
            return output_path

        except Exception as e:
            logger.exception("Image to PDF conversion failed:")
            raise ConversionError(
                "Failed to convert image to PDF. "
                "The image may be corrupted or in an unsupported format."
            ) from e

    @staticmethod
    async def word_to_pdf(input_path: Path, output_path: Path) -> Path:
        """Convert a Word document to PDF format.

        Reads formatting from python-docx (bold, italic, font size,
        headings, bullet lists) and renders a styled PDF using ReportLab.
        Not pixel-perfect compared to LibreOffice, but preserves the
        document's visual hierarchy and formatting.

        Args:
            input_path: Path to the source Word document.
            output_path: Where to save the generated PDF.

        Returns:
            Path to the created PDF file.

        Raises:
            ConversionError: If the Word document cannot be read or PDF creation fails.
        """
        try:
            from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem,
            )
            from reportlab.lib.units import inch

            doc = Document(str(input_path))

            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                leftMargin=72,
                rightMargin=72,
                topMargin=72,
                bottomMargin=72,
            )

            styles = getSampleStyleSheet()

            # Custom styles for different Word paragraph types
            styles.add(ParagraphStyle(
                "DocHeading1",
                parent=styles["Heading1"],
                fontSize=18,
                spaceAfter=12,
                spaceBefore=16,
            ))
            styles.add(ParagraphStyle(
                "DocHeading2",
                parent=styles["Heading2"],
                fontSize=14,
                spaceAfter=8,
                spaceBefore=12,
            ))
            styles.add(ParagraphStyle(
                "DocBody",
                parent=styles["Normal"],
                fontSize=11,
                leading=14,
                spaceAfter=4,
            ))
            styles.add(ParagraphStyle(
                "DocBodyCenter",
                parent=styles["Normal"],
                fontSize=11,
                leading=14,
                alignment=TA_CENTER,
                spaceAfter=4,
            ))

            story: list = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    story.append(Spacer(1, 6))
                    continue

                # Determine style from Word paragraph style name
                style_name = paragraph.style.name if paragraph.style else ""

                if "Heading 1" in style_name or "Title" in style_name:
                    style = styles["DocHeading1"]
                elif "Heading 2" in style_name:
                    style = styles["DocHeading2"]
                elif "Heading" in style_name:
                    style = styles["DocHeading2"]
                else:
                    # Check alignment
                    from docx.enum.text import WD_ALIGN_PARAGRAPH as WDA
                    if paragraph.alignment == WDA.CENTER:
                        style = styles["DocBodyCenter"]
                    else:
                        style = styles["DocBody"]

                # Build rich text with inline formatting from runs
                rich_parts: list[str] = []
                for run in paragraph.runs:
                    run_text = run.text
                    if not run_text:
                        continue

                    # Escape XML special chars for ReportLab
                    run_text = (
                        run_text.replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;")
                    )

                    if run.bold and run.italic:
                        rich_parts.append(f"<b><i>{run_text}</i></b>")
                    elif run.bold:
                        rich_parts.append(f"<b>{run_text}</b>")
                    elif run.italic:
                        rich_parts.append(f"<i>{run_text}</i>")
                    else:
                        rich_parts.append(run_text)

                rich_text = "".join(rich_parts) if rich_parts else text

                # Handle list/bullet paragraphs
                if "List" in style_name or "Bullet" in style_name:
                    bullet_para = Paragraph(rich_text, styles["DocBody"])
                    story.append(
                        ListFlowable(
                            [ListItem(bullet_para)],
                            bulletType="bullet",
                            start=None,
                        )
                    )
                else:
                    story.append(Paragraph(rich_text, style))

            pdf_doc.build(story)

            logger.info(
                "Word to PDF conversion complete: %s",
                output_path.name,
            )
            return output_path

        except Exception as e:
            logger.exception("Word to PDF conversion failed:")
            raise ConversionError(
                "Failed to convert Word document to PDF. "
                "The file may be corrupted or in an unsupported format."
            ) from e

    @staticmethod
    async def pdf_to_excel(input_path: Path, output_path: Path) -> Path:
        """Convert a PDF document to Excel (.xlsx) format.

        Extracts tabular data using PyMuPDF's built-in table detection.
        Falls back to placing each text line in column A when no tables
        are found on a page.

        Args:
            input_path: Path to the source PDF file.
            output_path: Where to save the generated Excel workbook.

        Returns:
            Path to the created Excel workbook.

        Raises:
            ConversionError: If the PDF cannot be read or Excel creation fails.
        """
        try:
            with fitz.open(str(input_path)) as doc:
                if len(doc) > 300:
                    raise ConversionError(
                        "PDF has too many pages (max 300). "
                        "Please split the document and try again."
                    )

                wb = Workbook()
                sheets_created = 0

                for page_num in range(len(doc)):
                    page = doc[page_num]
                    ws = wb.create_sheet(title=f"Page {page_num + 1}")
                    sheets_created += 1
                    row_cursor = 1

                    # Try built-in table detection (PyMuPDF 1.23+)
                    tables = page.find_tables()
                    if tables and len(tables.tables) > 0:
                        for table in tables.tables:
                            for row_data in table.extract():
                                for col_idx, cell_value in enumerate(row_data, start=1):
                                    ws.cell(
                                        row=row_cursor,
                                        column=col_idx,
                                        value=cell_value if cell_value else "",
                                    )
                                row_cursor += 1
                            # Blank row between tables
                            row_cursor += 1
                    else:
                        # Fallback: detect column structure using text positions
                        blocks = page.get_text("dict")["blocks"]
                        rows: dict[float, list[tuple[float, str]]] = {}

                        for block in blocks:
                            if block.get("type") != 0:  # Skip image blocks
                                continue
                            for line in block.get("lines", []):
                                # Group lines by y-center (rounded) to form rows
                                y_center = round(
                                    (line["bbox"][1] + line["bbox"][3]) / 2, 0
                                )
                                line_text = " ".join(
                                    span["text"] for span in line["spans"]
                                ).strip()
                                x_pos = line["bbox"][0]
                                if line_text:
                                    if y_center not in rows:
                                        rows[y_center] = []
                                    rows[y_center].append((x_pos, line_text))

                        # Write rows sorted by y-position, columns by x-position
                        for y_pos in sorted(rows.keys()):
                            cols = sorted(rows[y_pos], key=lambda c: c[0])
                            for col_idx, (_, text) in enumerate(cols, start=1):
                                ws.cell(
                                    row=row_cursor,
                                    column=col_idx,
                                    value=text,
                                )
                            row_cursor += 1

                # Remove the default empty sheet created by Workbook()
                if sheets_created > 0 and "Sheet" in wb.sheetnames:
                    del wb["Sheet"]

                wb.save(str(output_path))

                logger.info(
                    "PDF to Excel conversion complete: %d pages -> %s",
                    len(doc),
                    output_path.name,
                )
                return output_path

        except ConversionError:
            raise
        except Exception as e:
            logger.exception("PDF to Excel conversion failed:")
            raise ConversionError(
                "Failed to convert PDF to Excel format. "
                "The file may be corrupted or contain unsupported content."
            ) from e

    @staticmethod
    async def pdf_to_powerpoint(input_path: Path, output_path: Path) -> Path:
        """Convert a PDF document to PowerPoint (.pptx) format.

        Renders each PDF page as a full-bleed background image on a slide
        and includes the extracted text as speaker notes for accessibility.

        Args:
            input_path: Path to the source PDF file.
            output_path: Where to save the generated PowerPoint presentation.

        Returns:
            Path to the created PowerPoint presentation.

        Raises:
            ConversionError: If the PDF cannot be read or presentation creation fails.
        """
        try:
            with fitz.open(str(input_path)) as doc:
                if len(doc) > 300:
                    raise ConversionError(
                        "PDF has too many pages (max 300). "
                        "Please split the document and try again."
                    )

                prs = Presentation()
                # Set 16:9 slide dimensions (13.333 x 7.5 inches)
                prs.slide_width = Emu(int(13.333 * 914400))
                prs.slide_height = Emu(int(7.5 * 914400))

                # Use the blank slide layout (index 6 is typically blank)
                blank_layout = prs.slide_layouts[6]

                tmp_dir = output_path.parent

                for page_num in range(len(doc)):
                    page = doc[page_num]

                    # Render page as PNG at 150 DPI
                    zoom = 150 / 72
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    img_path = tmp_dir / f"slide_{page_num + 1}.png"
                    pix.save(str(img_path))

                    # Add blank slide
                    slide = prs.slides.add_slide(blank_layout)

                    # Add image as full-bleed background
                    slide.shapes.add_picture(
                        str(img_path),
                        left=Emu(0),
                        top=Emu(0),
                        width=prs.slide_width,
                        height=prs.slide_height,
                    )

                    # Extract text and add as speaker notes
                    text = page.get_text("text").strip()
                    if text:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = text

                prs.save(str(output_path))

                logger.info(
                    "PDF to PowerPoint conversion complete: %d pages -> %s",
                    len(doc),
                    output_path.name,
                )
                return output_path

        except ConversionError:
            raise
        except Exception as e:
            logger.exception("PDF to PowerPoint conversion failed:")
            raise ConversionError(
                "Failed to convert PDF to PowerPoint format. "
                "The file may be corrupted or contain unsupported content."
            ) from e

    @staticmethod
    async def excel_to_pdf(input_path: Path, output_path: Path) -> Path:
        """Convert an Excel workbook to PDF format.

        Reads each worksheet with openpyxl and renders tables using
        ReportLab's Platypus Table flowable.

        Args:
            input_path: Path to the source .xlsx file.
            output_path: Where to save the generated PDF.

        Returns:
            Path to the created PDF file.
        """
        try:
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Spacer, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors

            from openpyxl import load_workbook
            wb = load_workbook(str(input_path), data_only=True)

            pdf_doc = SimpleDocTemplate(
                str(output_path), pagesize=letter,
                leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36,
            )
            styles = getSampleStyleSheet()
            story: list = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                story.append(Paragraph(f"<b>{sheet_name}</b>", styles["Heading2"]))

                data: list[list[str]] = []
                for row in ws.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else "" for cell in row])

                if not data:
                    story.append(Paragraph("(empty sheet)", styles["Normal"]))
                    continue

                # Limit column widths to fit on page
                num_cols = max(len(r) for r in data) if data else 1
                col_width = min(120, (7.5 * 72) / num_cols)

                table = Table(data, colWidths=[col_width] * num_cols)
                table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.Color(0.95, 0.95, 0.95)]),
                ]))
                story.append(table)
                story.append(Spacer(1, 20))

            pdf_doc.build(story)
            logger.info("Excel to PDF conversion complete: %s", output_path.name)
            return output_path

        except Exception as e:
            logger.exception("Excel to PDF conversion failed:")
            raise ConversionError(
                "Failed to convert Excel to PDF. "
                "The file may be corrupted or in an unsupported format."
            ) from e

    @staticmethod
    async def pptx_to_pdf(input_path: Path, output_path: Path) -> Path:
        """Convert a PowerPoint presentation to PDF format.

        Renders each slide as an image by extracting text and shapes,
        then assembles into a PDF. For best fidelity, each slide is
        rendered as a descriptive page with title and content.

        Args:
            input_path: Path to the source .pptx file.
            output_path: Where to save the generated PDF.

        Returns:
            Path to the created PDF file.
        """
        try:
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.enums import TA_CENTER

            prs = Presentation(str(input_path))

            pdf_doc = SimpleDocTemplate(
                str(output_path), pagesize=letter,
                leftMargin=72, rightMargin=72, topMargin=72, bottomMargin=72,
            )
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(
                "SlideTitle", parent=styles["Heading1"],
                fontSize=20, alignment=TA_CENTER, spaceAfter=20,
            ))
            story: list = []

            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract title
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text

                # Escape XML special chars in title (ReportLab Paragraph uses XML markup)
                safe_title = (
                    title.replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                ) if title else ""
                story.append(Paragraph(
                    f"<b>Slide {slide_num}</b>" + (f": {safe_title}" if safe_title else ""),
                    styles["SlideTitle"],
                ))

                # Extract all text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and text != title:
                                # Escape XML
                                text = (
                                    text.replace("&", "&amp;")
                                    .replace("<", "&lt;")
                                    .replace(">", "&gt;")
                                )
                                story.append(Paragraph(text, styles["Normal"]))

                # Add page break between slides
                if slide_num < len(prs.slides):
                    story.append(PageBreak())

            pdf_doc.build(story)
            logger.info(
                "PPTX to PDF conversion complete: %d slides -> %s",
                len(prs.slides), output_path.name,
            )
            return output_path

        except Exception as e:
            logger.exception("PPTX to PDF conversion failed:")
            raise ConversionError(
                "Failed to convert PowerPoint to PDF. "
                "The file may be corrupted or in an unsupported format."
            ) from e

    @staticmethod
    async def html_to_pdf(input_path: Path, output_path: Path) -> Path:
        """Convert an HTML file to PDF format.

        Parses the HTML content and renders it to PDF using ReportLab's
        Paragraph flowable which supports basic HTML tags (b, i, br, etc.).

        Args:
            input_path: Path to the source .html file.
            output_path: Where to save the generated PDF.

        Returns:
            Path to the created PDF file.
        """
        try:
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            import re as re_mod

            html_content = input_path.read_text(encoding="utf-8")

            # Strip <script> and <style> tags
            html_content = re_mod.sub(
                r"<(script|style)[^>]*>.*?</\1>", "", html_content,
                flags=re_mod.DOTALL | re_mod.IGNORECASE,
            )

            # Extract title if present
            title_match = re_mod.search(
                r"<title[^>]*>(.*?)</title>",
                html_content, re_mod.IGNORECASE | re_mod.DOTALL,
            )

            # Extract body content
            body_match = re_mod.search(
                r"<body[^>]*>(.*?)</body>",
                html_content, re_mod.IGNORECASE | re_mod.DOTALL,
            )
            body = body_match.group(1) if body_match else html_content

            # Convert headings to bold
            for h in ("h1", "h2", "h3", "h4", "h5", "h6"):
                body = re_mod.sub(
                    rf"<{h}[^>]*>(.*?)</{h}>",
                    r"<b>\1</b><br/>",
                    body,
                    flags=re_mod.IGNORECASE | re_mod.DOTALL,
                )

            # Convert strong/em to b/i
            body = re_mod.sub(r"<strong[^>]*>", "<b>", body, flags=re_mod.IGNORECASE)
            body = re_mod.sub(r"</strong>", "</b>", body, flags=re_mod.IGNORECASE)
            body = re_mod.sub(r"<em[^>]*>", "<i>", body, flags=re_mod.IGNORECASE)
            body = re_mod.sub(r"</em>", "</i>", body, flags=re_mod.IGNORECASE)

            # Convert list items to bullet lines
            bullet_char = "\u2022"
            body = re_mod.sub(
                r"<li[^>]*>(.*?)</li>",
                lambda m: f"{bullet_char} {m.group(1)}<br/>",
                body,
                flags=re_mod.IGNORECASE | re_mod.DOTALL,
            )

            # Strip all remaining HTML tags except b, i, u, br
            cleaned = re_mod.sub(
                r"<(?!/?(?:b|i|u|br)\b)[^>]+>", " ", body
            )
            # Normalize whitespace
            cleaned = re_mod.sub(r"\s+", " ", cleaned).strip()

            pdf_doc = SimpleDocTemplate(
                str(output_path), pagesize=letter,
                leftMargin=72, rightMargin=72, topMargin=72, bottomMargin=72,
            )
            styles = getSampleStyleSheet()
            story: list = []

            if title_match:
                title = title_match.group(1).strip()
                story.append(Paragraph(f"<b>{title}</b>", styles["Title"]))
                story.append(Spacer(1, 12))

            # Split on <br/> and <p> to create separate paragraphs
            parts = re_mod.split(r"<br\s*/?>|</p>|<p[^>]*>", cleaned)
            for part in parts:
                text = part.strip()
                if text:
                    story.append(Paragraph(text, styles["Normal"]))

            if not story:
                story.append(Paragraph("(empty document)", styles["Normal"]))

            pdf_doc.build(story)
            logger.info("HTML to PDF conversion complete: %s", output_path.name)
            return output_path

        except Exception as e:
            logger.exception("HTML to PDF conversion failed:")
            raise ConversionError(
                "Failed to convert HTML to PDF. "
                "The file may be corrupted or in an unsupported format."
            ) from e
